"""
Generate group-meeting PPT (Chinese) for EC-PDNet multicase extension.

Output:
  paper/组会汇报_ECPDNet_全节点扩展.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
FIGURES = ROOT / "figures"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def add_title(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 248, 252)
    bg.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.8), Inches(1.5))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(18, 43, 78)

    s = slide.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(11.0), Inches(1.5))
    sp = s.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(58, 74, 95)


def add_bullets(prs: Presentation, title: str, bullets: list[str], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.8))
    hp = head.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(30)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(20, 48, 86)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(11.8), Inches(5.7))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(21)
        p.font.color.rgb = RGBColor(42, 56, 72)

    if note:
        n = slide.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(11.8), Inches(0.45))
        np = n.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(14)
        np.font.color.rgb = RGBColor(120, 120, 120)


def add_table_like(prs: Presentation, title: str, rows: list[list[str]], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.8))
    hp = head.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(30)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(20, 48, 86)

    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 1
    left = Inches(0.9)
    top = Inches(1.4)
    width = Inches(11.7)
    height = Inches(4.9)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    col_width = int(width / n_cols)
    for c in range(n_cols):
        table.columns[c].width = col_width

    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = rows[r][c]
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(16 if r > 0 else 17)
            para.font.bold = (r == 0)
            para.font.color.rgb = RGBColor(40, 50, 60)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(227, 238, 250)

    if note:
        n = slide.shapes.add_textbox(Inches(0.9), Inches(6.6), Inches(11.8), Inches(0.45))
        np = n.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(14)
        np.font.color.rgb = RGBColor(120, 120, 120)


def add_image(prs: Presentation, title: str, image: Path, caption: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.8))
    hp = head.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(30)
    hp.font.bold = True
    hp.font.color.rgb = RGBColor(20, 48, 86)

    if image.exists():
        slide.shapes.add_picture(str(image), Inches(0.9), Inches(1.2), width=Inches(11.8))
    else:
        box = slide.shapes.add_shape(1, Inches(0.9), Inches(1.2), Inches(11.8), Inches(5.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(205, 205, 205)

    cap = slide.shapes.add_textbox(Inches(0.95), Inches(6.85), Inches(11.7), Inches(0.4))
    cp = cap.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(90, 90, 90)


def main():
    summary = load_json(RESULTS / "ecpd_multicase_summary.json")
    case9 = load_json(RESULTS / "case9mod_ecpd_multicase_metrics.json")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title(
        prs,
        "组会汇报：EC-PDNet全节点系统扩展",
        "Peng Jiao | 目标：从安全域判别升级为物理一致的全状态代理",
    )

    add_bullets(
        prs,
        "1. 问题与目标",
        [
            "传统OPF/可行性扫描计算成本高，难以满足在线大规模点评估。",
            "仅输出可行/不可行不足以支撑调度决策，还需内部状态（P/Q/V/角度）。",
            "本阶段目标：构建并验证可迁移的EC-PDNet，覆盖WB2/WB5/case9/LMBM3。",
        ],
    )

    add_bullets(
        prs,
        "2. 方法亮点（中刊叙事）",
        [
            "坐标正确性：始终在发电机功率空间学习SSR，保持拓扑含义不变。",
            "能量闭合结构：由损耗头+闭合方程恢复关键状态，增强物理一致性。",
            "联合学习：分类头+状态头，兼顾边界判别与状态回归。",
            "可解释评估：分类指标 + 状态误差 + 闭合误差三维评价。",
        ],
        note="核心卖点：不是黑盒拟合，而是“结构化物理先验 + 数据驱动”的混合范式。",
    )

    rows = [["Case", "F1", "Acc", "State MAE", "Closure Abs Mean"]]
    for k in ["WB2", "WB5", "case9mod", "LMBM3_lf1p490", "LMBM3_lf1p500"]:
        if k not in summary:
            continue
        v = summary[k]
        def f(x):
            return "NaN" if (isinstance(x, float) and x != x) else f"{x:.4f}"
        rows.append([k, f(v["f1"]), f(v["acc"]), f(v["state_mae"]), f(v["closure_abs_mean"])])

    add_table_like(
        prs,
        "3. 全系统扩展结果（EC-PDNet）",
        rows,
        note="说明：WB5目前传统CSV仅含(PG1,PG5,loss)，状态指标维度与case9/LMBM3不完全一致。",
    )

    t = case9["test"]
    add_bullets(
        prs,
        "4. case9mod（全状态）重点结果",
        [
            f"分类：F1={t['classification']['f1']:.4f}, Acc={t['classification']['acc']:.4f}",
            f"状态：overall MAE={t['state']['overall_mae']:.4f}",
            f"分组：P={t['state']['mae_group'].get('p', float('nan')):.4f}, Q={t['state']['mae_group'].get('q', float('nan')):.4f}, V={t['state']['mae_group'].get('v', float('nan')):.4f}, theta={t['state']['mae_group'].get('theta', float('nan')):.4f}",
            f"闭合误差：mean={t['closure']['abs_mean']:.4f}, p95={t['closure']['abs_p95']:.4f}",
            "结论：结构化闭合约束显著提升物理一致性。",
        ],
    )

    add_image(
        prs,
        "5. 安全域拓扑复现示例（case9mod）",
        FIGURES / "case9mod_security_region.png",
        "在发电机功率空间下，模型保持对多连通安全域结构的复现能力。",
    )

    add_bullets(
        prs,
        "6. 论文改动与投稿策略",
        [
            "摘要改为单段，突出“坐标正确性+结构化闭合+全状态替代”主线。",
            "贡献压缩为3条：框架、方法创新、全基准验证与开源复现。",
            "实验部分新增EC-PDNet与多系统扩展结果，并给出CSV点对点对比文件。",
        ],
    )

    add_bullets(
        prs,
        "7. 下一步计划",
        [
            "补齐WB5传统状态导出（Q/V/theta），完善‘全状态’证据链。",
            "优化case9多系统统一训练策略，提升F1与状态精度平衡。",
            "增加推理时延对比（传统法 vs EC-PDNet）形成工程落地亮点。",
        ],
    )

    out = ROOT / "paper" / "组会汇报_ECPDNet_全节点扩展.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
