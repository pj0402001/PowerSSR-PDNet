"""
Generate Chinese project presentation PPT.

Output:
  paper/SSR_PDNet_全状态代理_汇报.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
FIGURES = ROOT / "figures"


def _load_metrics(name: str) -> dict:
    p = RESULTS / name
    return json.loads(p.read_text(encoding="utf-8"))


def _title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 48, 86)


def _note(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(6.65), Inches(11.7), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(120, 120, 120)


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 248, 252)
    bg.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.8), Inches(1.3))
    p = tbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(19, 46, 82)

    sbox = slide.shapes.add_textbox(Inches(0.85), Inches(2.5), Inches(11.0), Inches(1.8))
    p2 = sbox.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(54, 75, 96)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, title)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(11.8), Inches(5.7))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(21)
        p.font.color.rgb = RGBColor(42, 56, 72)
    if note:
        _note(slide, note)


def add_two_col_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_lines: list[str],
    right_title: str,
    right_lines: list[str],
    note: str | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, title)

    lbox = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(6.1), Inches(5.9))
    lbox.fill.solid()
    lbox.fill.fore_color.rgb = RGBColor(238, 245, 255)
    lbox.line.color.rgb = RGBColor(180, 200, 230)

    rbox = slide.shapes.add_shape(1, Inches(6.45), Inches(1.3), Inches(6.1), Inches(5.9))
    rbox.fill.solid()
    rbox.fill.fore_color.rgb = RGBColor(239, 250, 243)
    rbox.line.color.rgb = RGBColor(180, 220, 190)

    lt = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(5.6), Inches(0.5))
    lp = lt.text_frame.paragraphs[0]
    lp.text = left_title
    lp.font.size = Pt(20)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(24, 66, 120)

    lb = slide.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(5.6), Inches(4.8))
    ltf = lb.text_frame
    for i, line in enumerate(left_lines):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(45, 60, 80)

    rt = slide.shapes.add_textbox(Inches(6.7), Inches(1.55), Inches(5.6), Inches(0.5))
    rp = rt.text_frame.paragraphs[0]
    rp.text = right_title
    rp.font.size = Pt(20)
    rp.font.bold = True
    rp.font.color.rgb = RGBColor(22, 98, 42)

    rb = slide.shapes.add_textbox(Inches(6.7), Inches(2.05), Inches(5.6), Inches(4.8))
    rtf = rb.text_frame
    for i, line in enumerate(right_lines):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(45, 70, 52)

    if note:
        _note(slide, note)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, title)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.9), Inches(1.2), width=Inches(11.8))
    else:
        box = slide.shapes.add_shape(1, Inches(0.9), Inches(1.2), Inches(11.8), Inches(5.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(200, 200, 200)
        tx = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(11), Inches(0.6))
        p = tx.text_frame.paragraphs[0]
        p.text = f"图像缺失: {image_path.name}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(130, 130, 130)

    c = slide.shapes.add_textbox(Inches(0.95), Inches(6.85), Inches(11.7), Inches(0.4))
    cp = c.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(90, 90, 90)


def main():
    m_base = _load_metrics("case9mod_fullstate_pdnet_metrics.json")
    m_ecpd = _load_metrics("case9mod_fullstate_ecpd_metrics.json")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "基于物理闭合约束的电力系统安全域与全状态代理模型",
        "作者：Peng Jiao  |  目标：从“仅判别可行性”升级为“可行性+OPF状态联合替代”",
    )

    add_bullet_slide(
        prs,
        "1. 研究问题与痛点",
        [
            "传统 IPOPT 扫描准确但慢：每个点都要解非线性方程，在线应用成本高。",
            "仅做安全域分类不足以支撑调度：还需要 Q、V、角度等内部状态。",
            "投稿中刊需要“方法创新 + 物理解释 + 可复现实证”的完整闭环。",
        ],
        note="本工作核心：让模型既学边界拓扑，也学边界背后的物理状态。",
    )

    add_bullet_slide(
        prs,
        "2. 方法故事（创新主线）",
        [
            "提出 Energy-Closure PDNet（EC-PDNet）：把有功平衡闭合作为网络结构先验。",
            "创新A：不直接回归 P_G1，而是预测损耗 P_loss，再由闭合方程恢复 P_G1。",
            "创新B：状态分组头（Q/V/theta）+ 物理区间映射（V∈[0.9,1.1], Q有界）。",
            "创新C：联合目标（分类+状态）并引入单调先验（dP1/dP2, dP1/dP3 为负）。",
            "创新D：新增能量一致性指标，量化模型物理一致性。",
        ],
    )

    add_two_col_slide(
        prs,
        "3. 数学构想与关键公式",
        "结构化状态恢复",
        [
            "输入: x = [P_G2, P_G3]",
            "预测: P_loss = softplus(h_loss(z)) + c, c>0",
            "恢复: P_G1 = P_load + P_loss - P_G2 - P_G3",
            "Q/V/theta 由分组头输出并映射到物理可行区间。",
        ],
        "训练目标",
        [
            "L = L_cls + λ_state·L_state + λ_V·L_voltage + λ_mono·L_mono",
            "L_state: 仅在可行点监督（mask机制）",
            "L_mono: 约束 dP1/dP2, dP1/dP3 <= 0",
            "新增评估: E_close = |P1+P2+P3-P_load-P_loss_true|",
        ],
        note="把“可解释物理规律”嵌入网络结构与损失函数，避免纯黑盒拟合。",
    )

    add_two_col_slide(
        prs,
        "4. 实验设置（case9mod）",
        "数据与任务",
        [
            f"样本总数: {m_base['dataset']['n_total']}（可行 {m_base['dataset']['n_secure']} / 不可行 {m_base['dataset']['n_insecure']}）",
            "输入: (P_G2, P_G3)",
            "输出: 可行性 + 22维状态(P1,Q1~Q3,V1~V9,theta1~theta9)",
            "划分: 70/15/15",
        ],
        "对比模型",
        [
            "Baseline FullState-PDNet（直接状态回归）",
            "EC-PDNet（能量闭合结构化输出）",
            "指标: 分类Acc/F1 + 分组MAE + 能量一致性误差",
        ],
    )

    add_two_col_slide(
        prs,
        "5. 结果对比（测试集）",
        "Baseline FullState-PDNet",
        [
            f"F1 = {m_base['test']['classification']['f1']:.4f}",
            f"P1 MAE = {m_base['test']['state']['mae_group']['p_slack']:.4f} MW",
            f"Overall MAE = {m_base['test']['state']['overall_mae']:.4f}",
            f"Closure mean = {m_base['energy_consistency']['test']['closure_abs_mean_mw']:.4f} MW",
        ],
        "EC-PDNet（本次改进）",
        [
            f"F1 = {m_ecpd['test']['classification']['f1']:.4f}",
            f"P1 MAE = {m_ecpd['test']['state']['mae_group']['p_slack']:.4f} MW",
            f"Overall MAE = {m_ecpd['test']['state']['overall_mae']:.4f}",
            f"Closure mean = {m_ecpd['energy_consistency']['test']['closure_abs_mean_mw']:.4f} MW",
            "结论: P1误差和闭合误差显著下降，整体状态精度提升。",
        ],
    )

    add_bullet_slide(
        prs,
        "6. 关键结论（可写进摘要/结论）",
        [
            "EC-PDNet 将有功平衡内生化进网络，显著改善最难回归的 slack 有功 P_G1。",
            "P1 MAE 从约 0.86 MW 降到约 0.13 MW，能量闭合误差同步降低。",
            "分类性能保持高水平（F1仍约0.98），说明改进主要提升状态可信度。",
            "模型从“边界判别器”升级为“可替代传统点求解的全输出代理”。",
        ],
    )

    add_image_slide(
        prs,
        "7. 安全域拓扑复现（case9mod）",
        FIGURES / "case9mod_security_region.png",
        "模型保持对三连通安全域拓扑的重构能力（传统法与学习法边界一致性高）。",
    )

    add_bullet_slide(
        prs,
        "8. 工程落地流程",
        [
            "离线：传统法生成高质量样本，训练 EC-PDNet。",
            "在线：输入(P_G2,P_G3)直接输出可行概率与全状态变量。",
            "边界附近触发传统法复核（少量点），形成高效双轨推理。",
            "收益：绝大多数点无需迭代求解，显著加速在线评估。",
        ],
    )

    add_bullet_slide(
        prs,
        "9. 下一步工作",
        [
            "扩展 WB5/LMBM3 的全状态导出，形成跨系统统一证据。",
            "引入图结构编码（GNN）增强跨网络泛化。",
            "增加不确定性估计，构建可信度驱动的在线切换机制。",
            "补充复杂度分析：推理时延 vs 传统法时延。",
        ],
        note="建议将“结构化物理先验 + 可解释指标 + 开源复现”作为中刊创新主轴。",
    )

    out = ROOT / "paper" / "SSR_PDNet_全状态代理_汇报.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
